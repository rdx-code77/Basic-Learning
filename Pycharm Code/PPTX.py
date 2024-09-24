from pptx import Presentation
from pptx.util import Pt

# Create a presentation object
prs = Presentation()

# Slide titles
slide_titles = [
    "Dataset Identification and Access",
    "Source of Data",
    "Description of Dataset",
    "Accessing the Dataset",
    "Handling Missing Values",
    "Text Cleaning",
    "Stemming and Lemmatization"
]

# Content for each slide, split into bullet points
slide_contents = [
    ["In this step, we identify and access the dataset required for building the college website chatbot."],
    [
        "The data for this project is collected from the college website using web scraping techniques.",
        "Web scraping involves extracting data from web pages by parsing the HTML content of the website.",
        "For our project, we utilize the Beautiful Soup library in Python, which is a powerful tool for web scraping and parsing HTML/XML content.",
        "Beautiful Soup allows us to navigate and search the HTML structure of web pages, making it easier to extract specific data elements."
    ],
    [
        "After collecting the data from the college website using web scraping, we obtain a dataset that contains various types of information relevant to the chatbot.",
        "This dataset may include text content from different web pages on the college website, such as course descriptions, faculty profiles, admission information, etc.",
        "The format of the dataset may vary depending on the web scraping method used and the specific data elements extracted.",
        "Typically, the dataset is stored in a structured format such as CSV, JSON, or XML.",
        "The dataset may contain multiple columns or fields, each representing a different type of information."
    ],
    [
        "To access the dataset obtained from web scraping, we use Python code that utilizes the Beautiful Soup library to parse the HTML content of the college website.",
        "This code is executed within a Python script or Jupyter Notebook, allowing us to interactively access and explore the dataset.",
        "Once the data is collected and stored in a structured format, we can proceed to the next steps of the project."
    ],
    [
        "In this step, we address any missing or null values present in the dataset obtained from web scraping.",
        "Missing values can arise due to various reasons, such as incomplete data extraction or errors during web scraping.",
        "To handle missing values, we employ strategies such as imputation or deletion.",
        "Imputation involves replacing missing values with a calculated or estimated value based on the remaining data.",
        "We may choose to delete records with missing values if the proportion of missing values is small and does not significantly impact the overall dataset."
    ],
    [
        "The text data extracted from the college website may contain noise and irrelevant information that needs to be cleaned before further processing.",
        "Text cleaning involves several sub-steps such as Tokenization, Stopword Removal, Special Characters and Punctuation, and Lowercasing.",
        "We utilize tokenization techniques to split the text into individual tokens, which serve as the basic building blocks for further analysis.",
        "Stopwords are removed to improve processing efficiency using predefined lists or libraries such as NLTK.",
        "Special characters, punctuation marks, and symbols are removed or replaced to clean the text.",
        "All text is converted to lowercase to ensure consistency and avoid duplication of words during analysis."
    ],
    [
        "Stemming and lemmatization are techniques used to reduce words to their base or root forms, thereby reducing variations of words and improving text analysis.",
        "In stemming, words are chopped off to remove suffixes and prefixes.",
        "In lemmatization, words are reduced to their dictionary or canonical forms.",
        "We implement these techniques using libraries such as NLTK, which provides built-in functions for applying them to text data."
    ]
]

# Create slides
for title, content in zip(slide_titles, slide_contents):
    slide_layout = prs.slide_layouts[5]  # Using Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    content_placeholder = slide.placeholders[1]
    for bullet_point in content:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = bullet_point
        p.level = 0  # In case we need sub-levels, they can be adjusted here

# Save the presentation
pptx_path = '/mnt/data/Generated_Presentation.pptx'
prs.save(pptx_path)
